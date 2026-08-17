# Copyright (c) 2026 Trustedwear Tech Private Limited (https://citra-ai.com)
# Author: Rohit Kumar Chandan
# SPDX-License-Identifier: BUSL-1.1
#
# Licensed under the Business Source License 1.1. Non-production use is granted;
# production use requires a commercial licence until the Change Date, after
# which this file converts to Apache-2.0. See LICENSE at the repository root.

# ============================  Text Extractors  =============================
# Purpose: Text extraction utilities for multiple file types
# Features: Word, Excel, PowerPoint, Text, and enhanced PDF extraction
# Libraries: python-docx, pandas, python-pptx, PyMuPDF
# ----------------------------------------------------------------------------------------

import logging
import io
import os
import html as html_module
import traceback
import unicodedata
from pathlib import Path
from typing import Dict, List, Optional, Tuple
from dataclasses import dataclass

# ───────────────────────── Content Limits ──────────────────────────
# These limits are enforced during extraction to protect the system
MAX_PPT_SLIDES = int(os.getenv("MAX_PPT_SLIDES", 100))  # Max slides for PowerPoint
MAX_HTML_CHARS = int(os.getenv("MAX_HTML_CHARS", 200000))  # Max chars for HTML (200K)


# ─────────────── Unicode → ASCII normalization ───────────────
# Web pages frequently contain "smart" characters (curly quotes, em-dashes,
# non-breaking spaces, etc.) that look like garbled Unicode in plain-text
# consumers.  This map replaces them with ASCII equivalents.
_UNICODE_REPLACEMENTS = {
    # Quotation marks
    '\u2018': "'",   # left single curly quote
    '\u2019': "'",   # right single curly quote / apostrophe
    '\u201A': "'",   # single low-9 quotation mark
    '\u201B': "'",   # single high-reversed-9 quotation mark
    '\u201C': '"',   # left double curly quote
    '\u201D': '"',   # right double curly quote
    '\u201E': '"',   # double low-9 quotation mark
    '\u201F': '"',   # double high-reversed-9 quotation mark
    '\u2039': "'",   # single left-pointing angle quotation
    '\u203A': "'",   # single right-pointing angle quotation
    '\u00AB': '"',   # left-pointing double angle quotation «
    '\u00BB': '"',   # right-pointing double angle quotation »
    # Dashes / hyphens
    '\u2013': '-',   # en-dash
    '\u2014': '--',  # em-dash
    '\u2015': '--',  # horizontal bar
    '\u2212': '-',   # minus sign
    # Spaces
    '\u00A0': ' ',   # non-breaking space
    '\u2002': ' ',   # en space
    '\u2003': ' ',   # em space
    '\u2007': ' ',   # figure space
    '\u2008': ' ',   # punctuation space
    '\u2009': ' ',   # thin space
    '\u200A': ' ',   # hair space
    '\u202F': ' ',   # narrow no-break space
    '\u205F': ' ',   # medium mathematical space
    '\u3000': ' ',   # ideographic space
    # Zero-width / invisible characters (strip completely)
    '\u200B': '',    # zero-width space
    '\u200C': '',    # zero-width non-joiner
    '\u200D': '',    # zero-width joiner
    '\uFEFF': '',    # byte-order mark / zero-width no-break space
    '\u200E': '',    # left-to-right mark
    '\u200F': '',    # right-to-left mark
    # Misc punctuation
    '\u2026': '...',  # horizontal ellipsis
    '\u2022': '*',    # bullet
    '\u2023': '>',    # triangular bullet
    '\u2043': '-',    # hyphen bullet
    '\u25AA': '*',    # black small square (used as bullet)
    '\u25CF': '*',    # black circle (used as bullet)
    '\u00B7': '*',    # middle dot (used as bullet/separator)
    '\u2011': '-',    # non-breaking hyphen
    '\u00AD': '',     # soft hyphen
}

# Build a single translation table for str.translate() (fast, O(n) scan)
_UNICODE_TRANS_TABLE = str.maketrans({k: v for k, v in _UNICODE_REPLACEMENTS.items()})


def _normalize_unicode_to_ascii(text: str) -> str:
    """Replace common Unicode characters with ASCII equivalents.
    
    Applies NFKC normalization first (e.g. ﬁ→fi, ²→2), then maps
    smart quotes, special dashes/spaces, zero-width chars, etc. to ASCII.
    """
    # NFKC normalisation handles ligatures (ﬁ→fi), superscripts (²→2), etc.
    text = unicodedata.normalize('NFKC', text)
    # Fast character-level translation
    text = text.translate(_UNICODE_TRANS_TABLE)
    return text


# File type specific imports
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    logging.warning("python-docx not available. Word document support disabled.")

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    PANDAS_AVAILABLE = False
    logging.warning("pandas not available. Excel support disabled.")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logging.warning("python-pptx not available. PowerPoint support disabled.")

try:
    import fitz  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False
    logging.warning("PyMuPDF not available. Enhanced PDF support disabled.")

try:
    from googleapiclient.discovery import build
    from google.auth.transport.requests import Request
    from google_auth_oauthlib.flow import InstalledAppFlow
    from google.auth import default
    import json
    GOOGLE_API_AVAILABLE = True
except ImportError:
    GOOGLE_API_AVAILABLE = False
    logging.debug("Google API libraries not available. Direct Google Docs API access disabled (exported files still supported).")

# =========================== Source Information Tracking ===========================

class SourceInfo:
    """Class to track source information for extracted text"""
    def __init__(self, file_type: str, filename: str):
        self.file_type = file_type
        self.filename = filename
        self.sections = []
    
    def add_section(self, text: str, page_num: Optional[int] = None, 
                   section_type: str = "content", metadata: Optional[Dict] = None):
        """Add a text section with source information"""
        section = {
            "text": text,
            "page_number": page_num,
            "section_type": section_type,
            "metadata": metadata or {}
        }
        self.sections.append(section)
    
    def get_combined_text(self) -> str:
        """Get all text combined with source markers"""
        combined_parts = []
        
        for section in self.sections:
            text = section["text"].strip()
            if not text:
                continue
                
            # Add source markers for better citation
            if section["page_number"]:
                combined_parts.append(f"=== PAGE {section['page_number']} ===")
            elif section["section_type"] != "content":
                combined_parts.append(f"=== {section['section_type'].upper()} ===")
            
            combined_parts.append(text)
            combined_parts.append("")  # Add spacing between sections
        
        return "\n".join(combined_parts)
    
    def get_metadata_summary(self) -> Dict:
        """Get summary metadata for the extracted content"""
        total_pages = len([s for s in self.sections if s.get("page_number")])
        total_sections = len(self.sections)
        
        return {
            "file_type": self.file_type,
            "filename": self.filename,
            "total_pages": total_pages,
            "total_sections": total_sections,
            "extraction_method": f"{self.file_type}_direct",
            "source_preservation": True
        }


@dataclass
class DocAnalysis:
    """Lightweight analysis results for Word documents"""
    total_paragraphs: int
    tables: int
    images: int
    estimated_pages: int
    recommended_strategy: str
    processing_estimate: float
    has_mixed_content: bool

# =========================== Word Document Extraction ===========================

def extract_text_from_docx(file_content: bytes, filename: str) -> Tuple[str, Dict]:
    """
    Extract text from Word document (.docx) with source preservation
    """
    if not DOCX_AVAILABLE:
        raise ValueError("python-docx library not available. Cannot process Word documents.")
    
    try:
        logging.info(f"[DOCX] Starting text extraction from {filename}")
        
        # Create source info tracker
        source_info = SourceInfo("docx", filename)
        
        # Load document from bytes
        doc_stream = io.BytesIO(file_content)
        doc = Document(doc_stream)
        
        # Extract paragraphs with source tracking
        for i, paragraph in enumerate(doc.paragraphs):
            if paragraph.text.strip():
                source_info.add_section(
                    text=paragraph.text,
                    section_type="paragraph",
                    metadata={"paragraph_index": i}
                )
        
        # Extract text from tables
        table_count = 0
        for table in doc.tables:
            table_count += 1
            table_text = []
            
    # Return combined text and metadata as before
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    cell_text = cell.text.strip()
                    if cell_text:
                        row_text.append(cell_text)
                
                if row_text:
                    table_text.append(" | ".join(row_text))
            
            if table_text:
                source_info.add_section(
                    text="\n".join(table_text),
                    section_type=f"table_{table_count}",
                    metadata={"table_index": table_count}
                )
        
        combined_text = source_info.get_combined_text()
        metadata = source_info.get_metadata_summary()
        
        # Calculate page count for Word documents (approximation based on word count)
        word_count = len(combined_text.split())
        estimated_pages = max(1, round(word_count / 450))  # ~450 words per page typical
        metadata['total_pages'] = estimated_pages
        
        logging.info(f"[DOCX] Successfully extracted text from {filename}")
        logging.info(f"[DOCX] Found {metadata['total_sections']} sections, estimated {estimated_pages} pages")
        
        return combined_text, metadata
        
    except Exception as e:
        logging.error(f"[DOCX] Error extracting text from {filename}: {str(e)}")
        logging.error(f"[DOCX] Traceback: {traceback.format_exc()}")
        raise ValueError(f"Failed to extract text from Word document: {str(e)}")

# =========================== Excel Extraction ===========================

def extract_text_from_excel(file_content: bytes, filename: str) -> Tuple[str, Dict]:
    """
    Extract text from Excel files (.xlsx, .xls) with sheet and cell preservation
    """
    if not PANDAS_AVAILABLE:
        raise ValueError("pandas library not available. Cannot process Excel files.")
    
    try:
        logging.info(f"[EXCEL] Starting text extraction from {filename}")
        
        # Create source info tracker
        source_info = SourceInfo("excel", filename)
        
        # Load Excel file from bytes
        excel_stream = io.BytesIO(file_content)
        
        # Read all sheets
        excel_file = pd.ExcelFile(excel_stream)
        
        for sheet_name in excel_file.sheet_names:
            try:
                # Read sheet data
                df = pd.read_excel(excel_file, sheet_name=sheet_name)
                
                # Convert to string and clean
                sheet_text_parts = []
                
                # Add header information
                if not df.empty and not df.columns.empty:
                    headers = [str(col) for col in df.columns if str(col) != 'Unnamed: 0']
                    if headers:
                        sheet_text_parts.append("HEADERS: " + " | ".join(headers))
                
                # Add row data
                for index, row in df.iterrows():
                    row_values = []
                    for value in row:
                        if pd.notna(value) and str(value).strip():
                            row_values.append(str(value).strip())
                    
                    if row_values:
                        sheet_text_parts.append(" | ".join(row_values))
                
                if sheet_text_parts:
                    sheet_text = "\n".join(sheet_text_parts)
                    source_info.add_section(
                        text=sheet_text,
                        section_type=f"sheet",
                        metadata={
                            "sheet_name": sheet_name,
                            "rows": len(df),
                            "columns": len(df.columns)
                        }
                    )
                
            except Exception as sheet_error:
                logging.warning(f"[EXCEL] Error reading sheet '{sheet_name}': {str(sheet_error)}")
                continue
        
        combined_text = source_info.get_combined_text()
        metadata = source_info.get_metadata_summary()
        
        # Add Excel-specific metadata
        metadata.update({
            "total_sheets": len(excel_file.sheet_names),
            "sheet_names": excel_file.sheet_names,
            "total_pages": 1  # Excel files typically counted as 1 page for document counting
        })
        
        logging.info(f"[EXCEL] Successfully extracted text from {filename}")
        logging.info(f"[EXCEL] Processed {len(excel_file.sheet_names)} sheets, counted as 1 page")
        
        return combined_text, metadata
        
    except Exception as e:
        logging.error(f"[EXCEL] Error extracting text from {filename}: {str(e)}")
        logging.error(f"[EXCEL] Traceback: {traceback.format_exc()}")
        raise ValueError(f"Failed to extract text from Excel file: {str(e)}")

def extract_excel_markdown(file_content: bytes, filename: str) -> Tuple[str, Dict]:
    """
    Extract text from Excel files (.xlsx, .xls) as Markdown tables for better LLM comprehension.
    """
    if not PANDAS_AVAILABLE:
        raise ValueError("pandas library not available. Cannot process Excel files.")
    
    try:
        logging.info(f"[EXCEL-MD] Starting Markdown extraction from {filename}")
        
        # Create source info tracker
        source_info = SourceInfo("excel_markdown", filename)
        
        # Load Excel file from bytes
        excel_stream = io.BytesIO(file_content)
        excel_file = pd.ExcelFile(excel_stream)
        
        total_rows = 0
        
        for sheet_name in excel_file.sheet_names:
            try:
                # Read sheet data
                df = pd.read_excel(excel_file, sheet_name=sheet_name)
                
                if df.empty:
                    continue
                    
                # Clean data: Replace NaNs with empty string
                df = df.fillna("")
                
                # Convert to Markdown
                # default implementation of to_markdown() requires tabulate
                try:
                    markdown_table = df.to_markdown(index=False)
                except ImportError:
                    # Fallback if tabulate is not installed
                    logging.warning("tabulate not installed, manual markdown conversion")
                    headers = list(df.columns)
                    header_row = "| " + " | ".join(map(str, headers)) + " |"
                    separator_row = "| " + " | ".join(["---"] * len(headers)) + " |"
                    data_rows = []
                    for _, row in df.iterrows():
                        data_rows.append("| " + " | ".join(map(str, row.values)) + " |")
                    markdown_table = "\n".join([header_row, separator_row] + data_rows)

                if markdown_table:
                    # Add Sheet Header
                    sheet_content = f"### Sheet: {sheet_name}\n\n{markdown_table}"
                    
                    source_info.add_section(
                        text=sheet_content,
                        section_type="sheet",
                        metadata={
                            "sheet_name": sheet_name,
                            "rows": len(df),
                            "columns": len(df.columns)
                        }
                    )
                    total_rows += len(df)
                
            except Exception as sheet_error:
                logging.warning(f"[EXCEL-MD] Error reading sheet '{sheet_name}': {str(sheet_error)}")
                continue
        
        combined_text = source_info.get_combined_text()
        metadata = source_info.get_metadata_summary()
        
        # Add Excel-specific metadata
        metadata.update({
            "total_sheets": len(excel_file.sheet_names),
            "sheet_names": excel_file.sheet_names,
            "total_rows_all_sheets": total_rows,
            "total_pages": max(1, round(total_rows / 50)) # Estimate 50 rows per page
        })
        
        logging.info(f"[EXCEL-MD] Successfully extracted Markdown from {filename}")
        
        return combined_text, metadata
        
    except Exception as e:
        logging.error(f"[EXCEL-MD] Error extracting text from {filename}: {str(e)}")
        raise ValueError(f"Failed to extract Markdown from Excel file: {str(e)}")

# =========================== CSV Extraction ===========================

def extract_csv_markdown(file_content: bytes, filename: str) -> Tuple[str, Dict]:
    """
    Extract text from CSV files as Markdown tables for better LLM comprehension.
    Follows the same pattern as extract_excel_markdown for consistency.
    """
    if not PANDAS_AVAILABLE:
        raise ValueError("pandas library not available. Cannot process CSV files.")
    
    try:
        logging.info(f"[CSV-MD] Starting Markdown extraction from {filename}")
        
        # Create source info tracker
        source_info = SourceInfo("csv_markdown", filename)
        
        # Load CSV file from bytes
        csv_stream = io.BytesIO(file_content)
        df = pd.read_csv(csv_stream)
        
        total_rows = 0
        
        if not df.empty:
            # Clean data: Replace NaNs with empty string
            df = df.fillna("")
            
            # Convert to Markdown
            try:
                markdown_table = df.to_markdown(index=False)
            except ImportError:
                # Fallback if tabulate is not installed
                logging.warning("tabulate not installed, manual markdown conversion")
                headers = list(df.columns)
                header_row = "| " + " | ".join(map(str, headers)) + " |"
                separator_row = "| " + " | ".join(["---"] * len(headers)) + " |"
                data_rows = []
                for _, row in df.iterrows():
                    data_rows.append("| " + " | ".join(map(str, row.values)) + " |")
                markdown_table = "\n".join([header_row, separator_row] + data_rows)
            
            if markdown_table:
                sheet_content = f"### Data\n\n{markdown_table}"
                
                source_info.add_section(
                    text=sheet_content,
                    section_type="data",
                    metadata={
                        "rows": len(df),
                        "columns": len(df.columns)
                    }
                )
                total_rows = len(df)
        
        combined_text = source_info.get_combined_text()
        metadata = source_info.get_metadata_summary()
        
        # Add CSV-specific metadata
        metadata.update({
            "total_rows": total_rows,
            "total_rows_all_sheets": total_rows,
            "total_pages": max(1, round(total_rows / 50))  # Estimate 50 rows per page
        })
        
        logging.info(f"[CSV-MD] Successfully extracted Markdown from {filename}: {total_rows} rows")
        
        return combined_text, metadata
        
    except Exception as e:
        logging.error(f"[CSV-MD] Error extracting text from {filename}: {str(e)}")
        raise ValueError(f"Failed to extract Markdown from CSV file: {str(e)}")

# =========================== PowerPoint Extraction ===========================

def extract_text_from_pptx(file_content: bytes, filename: str) -> Tuple[str, Dict]:
    """
    Extract text from PowerPoint presentations (.pptx) with slide preservation
    """
    if not PPTX_AVAILABLE:
        raise ValueError("python-pptx library not available. Cannot process PowerPoint files.")
    
    try:
        logging.info(f"[PPTX] Starting text extraction from {filename}")
        
        # Create source info tracker
        source_info = SourceInfo("pptx", filename)
        
        # Load presentation from bytes
        pptx_stream = io.BytesIO(file_content)
        prs = Presentation(pptx_stream)
        
        # ========== PPT SLIDE LIMIT CHECK ==========
        total_slides = len(prs.slides)
        if total_slides > MAX_PPT_SLIDES:
            logging.warning(f"[PPTX] ⚠️ PowerPoint exceeds slide limit: {total_slides} slides (max: {MAX_PPT_SLIDES})")
            raise ValueError(f"PowerPoint exceeds maximum slide limit of {MAX_PPT_SLIDES} slides. Your file has {total_slides} slides.")
        # ========== END SLIDE LIMIT CHECK ==========
        
        slide_count = 0
        for slide in prs.slides:
            slide_count += 1
            slide_text_parts = []
            
            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text_parts.append(shape.text.strip())
                
                # Extract text from tables
                if hasattr(shape, "table"):
                    table_text = []
                    for row in shape.table.rows:
                        row_cells = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_cells.append(cell.text.strip())
                        if row_cells:
                            table_text.append(" | ".join(row_cells))
                    
                    if table_text:
                        slide_text_parts.append("TABLE DATA:\n" + "\n".join(table_text))
            
            if slide_text_parts:
                slide_text = "\n\n".join(slide_text_parts)
                source_info.add_section(
                    text=slide_text,
                    page_num=slide_count,
                    section_type="slide",
                    metadata={"slide_number": slide_count}
                )
        
        combined_text = source_info.get_combined_text()
        metadata = source_info.get_metadata_summary()
        
        # Add PowerPoint-specific metadata
        metadata.update({
            "total_slides": slide_count,
            "total_pages": slide_count  # Each slide counts as a page
        })
        
        logging.info(f"[PPTX] Successfully extracted text from {filename}")
        logging.info(f"[PPTX] Processed {slide_count} slides ({slide_count} pages)")
        
        return combined_text, metadata
        
    except Exception as e:
        logging.error(f"[PPTX] Error extracting text from {filename}: {str(e)}")
        logging.error(f"[PPTX] Traceback: {traceback.format_exc()}")
        raise ValueError(f"Failed to extract text from PowerPoint file: {str(e)}")

# =========================== Text File Extraction ===========================

def extract_text_from_txt(file_content: bytes, filename: str) -> Tuple[str, Dict]:
    """
    Extract text from plain text files (.txt, .md) with encoding detection
    """
    try:
        logging.info(f"[TXT] Starting text extraction from {filename}")
        
        # Try different encodings
        encodings = ['utf-8', 'utf-8-sig', 'latin1', 'cp1252']
        text_content = None
        used_encoding = None
        
        for encoding in encodings:
            try:
                text_content = file_content.decode(encoding)
                used_encoding = encoding
                break
            except UnicodeDecodeError:
                continue
        
        if text_content is None:
            # Last resort: decode with errors='replace'
            text_content = file_content.decode('utf-8', errors='replace')
            used_encoding = 'utf-8-with-replacements'
        
        # Create source info tracker
        source_info = SourceInfo("txt", filename)
        
        # For text files, we can split by paragraphs or sections
        lines = text_content.split('\n')
        current_paragraph = []
        paragraph_count = 0
        
        for line in lines:
            line = line.strip()
            
            if line:  # Non-empty line
                current_paragraph.append(line)
            else:  # Empty line - end of paragraph
                if current_paragraph:
                    paragraph_count += 1
                    paragraph_text = "\n".join(current_paragraph)
                    source_info.add_section(
                        text=paragraph_text,
                        section_type="paragraph",
                        metadata={
                            "paragraph_number": paragraph_count,
                            "line_count": len(current_paragraph)
                        }
                    )
                    current_paragraph = []
        
        # Add final paragraph if exists
        if current_paragraph:
            paragraph_count += 1
            paragraph_text = "\n".join(current_paragraph)
            source_info.add_section(
                text=paragraph_text,
                section_type="paragraph",
                metadata={
                    "paragraph_number": paragraph_count,
                    "line_count": len(current_paragraph)
                }
            )
        
        combined_text = source_info.get_combined_text()
        metadata = source_info.get_metadata_summary()
        
        # Add text-specific metadata
        word_count = len(combined_text.split())
        estimated_pages = max(1, round(word_count / 450))  # ~450 words per page typical
        metadata.update({
            "encoding_used": used_encoding,
            "total_lines": len(lines),
            "total_paragraphs": paragraph_count,
            "total_pages": estimated_pages
        })
        
        logging.info(f"[TXT] Successfully extracted text from {filename}")
        logging.info(f"[TXT] Used encoding: {used_encoding}, found {paragraph_count} paragraphs, estimated {estimated_pages} pages")
        
        return combined_text, metadata
        
    except Exception as e:
        logging.error(f"[TXT] Error extracting text from {filename}: {str(e)}")
        logging.error(f"[TXT] Traceback: {traceback.format_exc()}")
        raise ValueError(f"Failed to extract text from text file: {str(e)}")

# =========================== HTML Extraction ===========================

# Try to import BeautifulSoup for HTML parsing
try:
    from bs4 import BeautifulSoup
    BS4_AVAILABLE = True
    logging.info("✅ BeautifulSoup4 successfully imported for HTML parsing")
except ImportError:
    BS4_AVAILABLE = False
    logging.warning("beautifulsoup4 not available. HTML parsing support disabled.")

def extract_text_from_html(file_content: bytes, filename: str) -> Tuple[str, Dict]:
    """
    Extract text from HTML files with intelligent content parsing.
    Removes scripts, styles, navigation, and extracts meaningful content.
    Preserves document structure (headings, paragraphs, lists, tables).
    """
    if not BS4_AVAILABLE:
        raise ValueError("beautifulsoup4 library not available. Cannot process HTML files. Install with: pip install beautifulsoup4")
    
    try:
        logging.info(f"[HTML] Starting text extraction from {filename}")
        logging.info(f"[HTML] File size: {len(file_content)} bytes")
        
        # Create source info tracker
        source_info = SourceInfo("html", filename)
        
        # Try different encodings for HTML
        encodings = ['utf-8', 'utf-8-sig', 'latin1', 'cp1252', 'iso-8859-1']
        html_content = None
        used_encoding = None
        
        for encoding in encodings:
            try:
                html_content = file_content.decode(encoding)
                used_encoding = encoding
                logging.info(f"[HTML] Successfully decoded with encoding: {encoding}")
                break
            except UnicodeDecodeError:
                continue
        
        if html_content is None:
            raise ValueError(f"Failed to decode HTML file with any supported encoding: {encodings}")
        
        # Parse HTML with BeautifulSoup
        soup = BeautifulSoup(html_content, 'html.parser')
        logging.info(f"[HTML] Parsed HTML with BeautifulSoup, encoding: {used_encoding}")
        
        # Remove truly unwanted elements (scripts, chrome, boilerplate).
        # NOTE: 'aside' is intentionally NOT in this list — many sites (e.g. Wikipedia)
        # store valuable structured data (infoboxes, stats, coordinates) in <aside> tags.
        unwanted_tags = [
            'script', 'style', 'nav', 'footer', 'header',
            'noscript', 'iframe', 'form', 'button', 'input', 'select',
            'meta', 'link', 'comment'
        ]
        for tag in unwanted_tags:
            elements_removed = len(soup.find_all(tag))
            for element in soup.find_all(tag):
                element.decompose()
            if elements_removed > 0:
                logging.info(f"[HTML] Removed {elements_removed} <{tag}> elements")

        # ── Smart container selection ─────────────────────────────────────────
        # Try progressively broader selectors and pick the first one with ≥100 words.
        # Wikipedia Vector 2022 skin hides article text deep inside #mw-content-text;
        # a plain soup.find('main') returns only metadata/categories (≈287 words).
        _container_finders = [
            # MediaWiki / Wikipedia (Vector 2022 & older skins)
            lambda s: s.find('div', id='mw-content-text'),
            lambda s: s.find('div', id='bodyContent'),
            lambda s: s.find('div', class_='mw-parser-output'),
            # Common semantic / framework patterns
            lambda s: s.find('article'),
            lambda s: s.find('div', id='article-body'),
            lambda s: s.find('div', id='article-content'),
            lambda s: s.find('div', id='main-content'),
            lambda s: s.find('div', role='main'),
            lambda s: s.find('div', id='content'),
            lambda s: s.find('main'),
            lambda s: s.find('div', class_='content'),
            lambda s: s.find('div', class_='post-content'),
            lambda s: s.find('div', class_='article-body'),
            # Last resorts
            lambda s: s.body,
            lambda s: s,
        ]

        main_content = None
        for _finder in _container_finders:
            try:
                _candidate = _finder(soup)
                if _candidate is None:
                    continue
                _candidate_words = len(_candidate.get_text(separator=' ', strip=True).split())
                if _candidate_words >= 100:
                    main_content = _candidate
                    _tag = _candidate.name
                    _cid = _candidate.get('id', '')
                    _cls = _candidate.get('class', '')
                    logging.info(
                        f"[HTML] Selected content container: <{_tag}"
                        f"{' id=' + repr(_cid) if _cid else ''}"
                        f"{' class=' + repr(_cls) if _cls else ''}>"
                        f" ({_candidate_words} words)"
                    )
                    break
            except Exception:
                continue

        if main_content is None:
            main_content = soup.body or soup
            logging.warning("[HTML] No suitable content container found (≥100 words); falling back to <body>")
        else:
            logging.info(f"[HTML] Main content container: {main_content.name if hasattr(main_content, 'name') else 'None'}")
        
        # Get text length before ad removal
        text_before_cleanup = main_content.get_text(separator=' ', strip=True) if main_content else ""
        words_before = len(text_before_cleanup.split())
        logging.info(f"[HTML] Text before ad removal: {words_before} words")
        
        # Remove common ad/navigation elements (more conservative - exact class matches only)
        ad_selectors = [
            'nav[class*="nav"]', 'header[class*="header"]', 'footer[class*="footer"]',
            '[class*="cookie"]', '[class*="popup"]', '[class*="modal"]', 
            '[class*="advertisement"]', '[class*="banner"]'
        ]
        removed_ad_elements = 0
        for selector in ad_selectors:
            try:
                for element in main_content.select(selector):
                    element.decompose()
                    removed_ad_elements += 1
            except (ValueError, NotImplementedError):  # best-effort: skip selectors BeautifulSoup can't parse/support
                pass
        
        if removed_ad_elements > 0:
            logging.info(f"[HTML] Removed {removed_ad_elements} ad/navigation elements")
        
        # Simply extract all visible text from cleaned HTML
        extracted_text = main_content.get_text(separator=' ', strip=True) if main_content else ""
        
        # Unescape HTML entities (e.g. &amp; → &, &#8217; → ', &nbsp; → space)
        extracted_text = html_module.unescape(extracted_text)
        
        # Normalize Unicode → ASCII (smart quotes, em-dashes, non-breaking spaces, etc.)
        extracted_text = _normalize_unicode_to_ascii(extracted_text)
        
        # Clean up excessive whitespace
        extracted_text = ' '.join(extracted_text.split())
        
        word_count = len(extracted_text.split())
        logging.info(f"[HTML] Extracted {word_count} words from cleaned HTML")
        
        # ========== HTML CHARACTER LIMIT CHECK (TRUNCATE) ==========
        original_length = len(extracted_text)
        truncated = False
        if original_length > MAX_HTML_CHARS:
            logging.warning(f"[HTML] ⚠️ HTML content exceeds character limit: {original_length} chars (max: {MAX_HTML_CHARS}). Truncating...")
            extracted_text = extracted_text[:MAX_HTML_CHARS]
            word_count = len(extracted_text.split())
            truncated = True
            logging.info(f"[HTML] Truncated to {MAX_HTML_CHARS} chars ({word_count} words)")
        # ========== END CHARACTER LIMIT CHECK ==========
        
        # Calculate page estimate
        estimated_pages = max(1, round(word_count / 450))
        
        metadata = {
            "extraction_method": "beautifulsoup_get_text",
            "container": main_content.name if main_content and hasattr(main_content, 'name') else 'None',
            "encoding_used": used_encoding,
            "total_pages": estimated_pages,
            "word_count": word_count,
            "original_html_size": len(file_content),
            "extracted_text_size": len(extracted_text),
            "truncated": truncated,
            "original_text_length": original_length if truncated else len(extracted_text)
        }
        
        logging.info(f"[HTML] Successfully extracted text from {filename}")
        logging.info(f"[HTML] Extracted {word_count} words, estimated {estimated_pages} pages{' (truncated)' if truncated else ''}")
        
        return extracted_text, metadata
        
    except Exception as e:
        logging.error(f"[HTML] Error extracting text from {filename}: {str(e)}")
        logging.error(f"[HTML] Traceback: {traceback.format_exc()}")
        raise ValueError(f"Failed to extract text from HTML file: {str(e)}")


# =========================== JSON Extraction ===========================

def extract_text_from_json(file_content: bytes, filename: str) -> Tuple[str, Dict]:
    """
    Extract text from JSON files by flattening the structure and preserving context.
    Converts nested JSON into readable text format suitable for LLM context retrieval.
    """
    import json as json_module
    
    try:
        logging.info(f"[JSON] Starting text extraction from {filename}")
        
        # Create source info tracker
        source_info = SourceInfo("json", filename)
        
        # Try different encodings
        encodings = ['utf-8', 'utf-8-sig', 'latin1', 'cp1252']
        json_str = None
        used_encoding = None
        
        for encoding in encodings:
            try:
                json_str = file_content.decode(encoding)
                used_encoding = encoding
                break
            except UnicodeDecodeError:
                continue
        
        if json_str is None:
            json_str = file_content.decode('utf-8', errors='replace')
            used_encoding = 'utf-8-with-replacements'
        
        # Parse JSON to validate structure
        try:
            data = json_module.loads(json_str)
        except json_module.JSONDecodeError as e:
            raise ValueError(f"Invalid JSON format: {str(e)}")
        
        # Prettify JSON for LLM readability
        prettified_json = json_module.dumps(data, indent=2, ensure_ascii=False)
        
        word_count = len(prettified_json.split())
        estimated_pages = max(1, round(word_count / 450))
        
        logging.info(f"[JSON] Prettified JSON with {word_count} words")
        
        metadata = {
            "extraction_method": "json_prettified",
            "encoding_used": used_encoding,
            "total_pages": estimated_pages,
            "word_count": word_count,
            "original_json_size": len(file_content),
            "prettified_json_size": len(prettified_json)
        }
        
        logging.info(f"[JSON] Successfully prettified JSON from {filename}")
        
        return prettified_json, metadata
        
    except Exception as e:
        logging.error(f"[JSON] Error extracting text from {filename}: {str(e)}")
        logging.error(f"[JSON] Traceback: {traceback.format_exc()}")
        raise ValueError(f"Failed to extract text from JSON file: {str(e)}")


def analyze_docx_content(file_content: bytes, document_id: str) -> DocAnalysis:
    """
    Simple analysis for .docx files to estimate structure and whether enhanced processing is needed.
    Returns a DocAnalysis dataclass.
    """
    try:
        if not DOCX_AVAILABLE:
            # Conservative defaults when python-docx is not installed
            return DocAnalysis(
                total_paragraphs=0,
                tables=0,
                images=0,
                estimated_pages=1,
                recommended_strategy="text_primary",
                processing_estimate=1.0,
                has_mixed_content=False
            )

        doc = Document(io.BytesIO(file_content))

        paragraphs = [p for p in doc.paragraphs if p.text and p.text.strip()]
        total_paragraphs = len(paragraphs)
        table_count = len(doc.tables)

        # Count image relationships conservatively
        images = 0
        try:
            for rel in doc.part.rels.values():
                if hasattr(rel, 'reltype') and 'image' in str(rel.reltype).lower():
                    images += 1
        except Exception:
            images = 0

        # Estimate pages by words (approx. 450 words per page)
        total_words = sum(len(p.text.split()) for p in paragraphs)
        estimated_pages = max(1, int(total_words / 450) if total_words > 0 else 1)

        # Decide strategy: if images are frequent or many tables, prefer hybrid
        image_ratio = images / max(1, estimated_pages)
        if image_ratio > 0.5 or table_count >= max(1, estimated_pages // 2):
            strategy = "hybrid"
            processing_estimate = estimated_pages * 0.8
        else:
            strategy = "text_primary"
            processing_estimate = estimated_pages * 0.2

        has_mixed = (images > 0 and total_paragraphs > 0)

        return DocAnalysis(
            total_paragraphs=total_paragraphs,
            tables=table_count,
            images=images,
            estimated_pages=estimated_pages,
            recommended_strategy=strategy,
            processing_estimate=processing_estimate,
            has_mixed_content=has_mixed
        )

    except Exception as e:
        logging.warning(f"[DOCX_ANALYZE] Failed to analyze docx {document_id}: {e}")
        return DocAnalysis(
            total_paragraphs=0,
            tables=0,
            images=0,
            estimated_pages=1,
            recommended_strategy="text_primary",
            processing_estimate=1.0,
            has_mixed_content=False
        )

# =========================== Enhanced PDF Extraction ===========================

def extract_text_from_pdf_direct(file_content: bytes, filename: str) -> Tuple[str, Dict]:
    """
    Direct PDF text extraction using PyMuPDF (for text-only processing)
    This is used when Vision API is not needed for the new file types workflow
    """
    if not PYMUPDF_AVAILABLE:
        raise ValueError("PyMuPDF library not available. Cannot process PDF files.")
    
    try:
        logging.info(f"[PDF] Starting direct text extraction from {filename}")
        
        # Create source info tracker
        source_info = SourceInfo("pdf", filename)
        
        # Open PDF from bytes
        doc = fitz.open(stream=file_content, filetype="pdf")
        
        # Store page count before closing document
        total_pages = len(doc)
        
        for page_num in range(total_pages):
            page = doc.load_page(page_num)
            
            # Extract text from page
            page_text = page.get_text()
            
            if page_text.strip():
                source_info.add_section(
                    text=page_text,
                    page_num=page_num + 1,
                    section_type="page",
                    metadata={"page_number": page_num + 1}
                )
        
        doc.close()
        
        combined_text = source_info.get_combined_text()
        metadata = source_info.get_metadata_summary()
        
        # Add PDF-specific metadata
        metadata.update({
            "total_pages": total_pages,
            "extraction_type": "text_only"
        })
        
        logging.info(f"[PDF] Successfully extracted text from {filename}")
        logging.info(f"[PDF] Processed {total_pages} pages")
        
        return combined_text, metadata
        
    except Exception as e:
        logging.error(f"[PDF] Error extracting text from {filename}: {str(e)}")
        logging.error(f"[PDF] Traceback: {traceback.format_exc()}")
        raise ValueError(f"Failed to extract text from PDF file: {str(e)}")

# =========================== Google Docs Extraction ===========================

def extract_text_from_google_docs(file_id: str, filename: str, service_type: str = "docs") -> Tuple[str, Dict]:
    """
    Extract text from Google Docs, Sheets, or Slides using Google API
    Note: This requires proper Google API credentials and authentication
    """
    if not GOOGLE_API_AVAILABLE:
        raise ValueError("Google API libraries not available. Cannot process Google Docs files.")
    
    try:
        logging.info(f"[GOOGLE-{service_type.upper()}] Starting text extraction from {filename}")
        
        # Create source info tracker
        source_info = SourceInfo(f"google_{service_type}", filename)
        
        if service_type == "docs":
            return _extract_from_google_docs(file_id, filename, source_info)
        elif service_type == "sheets":
            return _extract_from_google_sheets(file_id, filename, source_info)
        elif service_type == "slides":
            return _extract_from_google_slides(file_id, filename, source_info)
        else:
            raise ValueError(f"Unsupported Google service type: {service_type}")
            
    except Exception as e:
        logging.error(f"[GOOGLE-{service_type.upper()}] Error extracting text from {filename}: {str(e)}")
        logging.error(f"[GOOGLE-{service_type.upper()}] Traceback: {traceback.format_exc()}")
        raise ValueError(f"Failed to extract text from Google {service_type}: {str(e)}")

def _get_google_service(service_name: str, version: str):
    """
    Get authenticated Google API service
    This function should be customized based on your authentication method
    """
    try:
        # Try to use default credentials (for service accounts or gcloud auth)
        credentials, project = default()
        service = build(service_name, version, credentials=credentials)
        return service
    except Exception as e:
        logging.error(f"Failed to authenticate with Google API: {str(e)}")
        raise ValueError("Google API authentication failed. Please configure credentials.")

def _extract_from_google_docs(file_id: str, filename: str, source_info: SourceInfo) -> Tuple[str, Dict]:
    """Extract text from Google Docs document"""
    service = _get_google_service('docs', 'v1')
    
    # Get document content
    document = service.documents().get(documentId=file_id).execute()
    
    # Extract text from document structure
    content = document.get('body', {}).get('content', [])
    
    for element in content:
        if 'paragraph' in element:
            paragraph = element['paragraph']
            paragraph_text = []
            
            for text_element in paragraph.get('elements', []):
                if 'textRun' in text_element:
                    text_content = text_element['textRun'].get('content', '')
                    if text_content.strip():
                        paragraph_text.append(text_content.strip())
            
            if paragraph_text:
                source_info.add_section(
                    text=' '.join(paragraph_text),
                    section_type="paragraph",
                    metadata={"source": "google_docs"}
                )
        
        elif 'table' in element:
            table = element['table']
            table_text = []
            
            for row in table.get('tableRows', []):
                row_cells = []
                for cell in row.get('tableCells', []):
                    cell_text = []
                    for cell_content in cell.get('content', []):
                        if 'paragraph' in cell_content:
                            for text_element in cell_content['paragraph'].get('elements', []):
                                if 'textRun' in text_element:
                                    text_content = text_element['textRun'].get('content', '')
                                    if text_content.strip():
                                        cell_text.append(text_content.strip())
                    
                    if cell_text:
                        row_cells.append(' '.join(cell_text))
                
                if row_cells:
                    table_text.append(' | '.join(row_cells))
            
            if table_text:
                source_info.add_section(
                    text='\n'.join(table_text),
                    section_type="table",
                    metadata={"source": "google_docs"}
                )
    
    combined_text = source_info.get_combined_text()
    metadata = source_info.get_metadata_summary()
    metadata.update({"google_docs_id": file_id})
    
    return combined_text, metadata

def _extract_from_google_sheets(file_id: str, filename: str, source_info: SourceInfo) -> Tuple[str, Dict]:
    """Extract text from Google Sheets document"""
    service = _get_google_service('sheets', 'v4')
    
    # Get spreadsheet metadata
    spreadsheet = service.spreadsheets().get(spreadsheetId=file_id).execute()
    sheets = spreadsheet.get('sheets', [])
    
    for sheet in sheets:
        sheet_properties = sheet.get('properties', {})
        sheet_title = sheet_properties.get('title', 'Unknown Sheet')
        
        # Get data from this sheet
        range_name = f"'{sheet_title}'"
        result = service.spreadsheets().values().get(
            spreadsheetId=file_id,
            range=range_name
        ).execute()
        
        values = result.get('values', [])
        if values:
            sheet_text = []
            for row in values:
                row_text = []
                for cell in row:
                    if str(cell).strip():
                        row_text.append(str(cell).strip())
                
                if row_text:
                    sheet_text.append(' | '.join(row_text))
            
            if sheet_text:
                source_info.add_section(
                    text='\n'.join(sheet_text),
                    section_type="sheet",
                    metadata={
                        "sheet_name": sheet_title,
                        "source": "google_sheets"
                    }
                )
    
    combined_text = source_info.get_combined_text()
    metadata = source_info.get_metadata_summary()
    metadata.update({
        "google_sheets_id": file_id,
        "total_sheets": len(sheets)
    })
    
    return combined_text, metadata

def _extract_from_google_slides(file_id: str, filename: str, source_info: SourceInfo) -> Tuple[str, Dict]:
    """Extract text from Google Slides presentation"""
    service = _get_google_service('slides', 'v1')
    
    # Get presentation content
    presentation = service.presentations().get(presentationId=file_id).execute()
    slides = presentation.get('slides', [])
    
    for slide_index, slide in enumerate(slides):
        slide_number = slide_index + 1
        slide_text = []
        
        for page_element in slide.get('pageElements', []):
            if 'shape' in page_element:
                shape = page_element['shape']
                if 'text' in shape:
                    text_elements = shape['text'].get('textElements', [])
                    for text_element in text_elements:
                        if 'textRun' in text_element:
                            text_content = text_element['textRun'].get('content', '')
                            if text_content.strip():
                                slide_text.append(text_content.strip())
            
            elif 'table' in page_element:
                table = page_element['table']
                table_text = []
                
                for row in table.get('tableRows', []):
                    row_cells = []
                    for cell in row.get('tableCells', []):
                        cell_text = []
                        if 'text' in cell:
                            for text_element in cell['text'].get('textElements', []):
                                if 'textRun' in text_element:
                                    text_content = text_element['textRun'].get('content', '')
                                    if text_content.strip():
                                        cell_text.append(text_content.strip())
                        
                        if cell_text:
                            row_cells.append(' '.join(cell_text))
                    
                    if row_cells:
                        table_text.append(' | '.join(row_cells))
                
                if table_text:
                    slide_text.append("TABLE:\n" + '\n'.join(table_text))
        
        if slide_text:
            source_info.add_section(
                text='\n\n'.join(slide_text),
                page_num=slide_number,
                section_type="slide",
                metadata={
                    "slide_number": slide_number,
                    "source": "google_slides"
                }
            )
    
    combined_text = source_info.get_combined_text()
    metadata = source_info.get_metadata_summary()
    metadata.update({
        "google_slides_id": file_id,
        "total_slides": len(slides)
    })
    
    return combined_text, metadata

# =========================== Google Docs File Upload Processing ===========================

def extract_text_from_google_file(file_content: bytes, filename: str, file_extension: str) -> Tuple[str, Dict]:
    """
    Process uploaded Google Docs files (exported as standard formats)
    Google Docs files can be exported as .docx, .xlsx, .pptx and processed normally
    """
    file_ext = file_extension.lower().lstrip('.')
    
    # Map Google file extensions to standard processing
    google_format_mapping = {
        'gdoc': 'docx',    # Google Docs exported as Word
        'gsheet': 'xlsx',  # Google Sheets exported as Excel  
        'gslides': 'pptx'  # Google Slides exported as PowerPoint
    }
    
    # If it's a Google format, map to standard format
    if file_ext in google_format_mapping:
        standard_format = google_format_mapping[file_ext]
        logging.info(f"[GOOGLE] Processing Google {file_ext} as {standard_format}")
        
        # Create modified metadata to indicate Google origin
        text, metadata = extract_text_by_file_type(file_content, filename, standard_format)
        metadata['original_format'] = f"google_{file_ext}"
        metadata['processing_method'] = f"google_export_as_{standard_format}"
        
        return text, metadata
    
    # If it's a standard format that could be from Google, process normally
    return extract_text_by_file_type(file_content, filename, file_extension)


# =========================== Legacy .doc Extraction ===========================
def extract_text_from_doc(file_content: bytes, filename: str) -> Tuple[str, Dict]:
    """
    Best-effort extraction for legacy .doc (binary Word) files.
    Tries pypandoc (requires pandoc) then textract. Returns (text, metadata).
    If neither tool is available, raises ValueError explaining the requirement.
    """
    import tempfile
    import os

    metadata = {"file_type": "doc", "filename": filename}

    # Save the bytes to a temporary .doc file
    with tempfile.NamedTemporaryFile(suffix='.doc', delete=False) as tmp:
        tmp.write(file_content)
        tmp_path = tmp.name

    try:
        # Try pypandoc first (requires pandoc installed on the system)
        try:
            import pypandoc
            try:
                text = pypandoc.convert_file(tmp_path, 'plain')
                # Add page count estimation for .doc files
                word_count = len(text.split())
                estimated_pages = max(1, round(word_count / 450))  # ~450 words per page
                metadata['total_pages'] = estimated_pages
                return text, metadata
            except Exception:
                # fall through to textract
                pass
        except Exception:
            # pypandoc not available
            pass

        # Try textract (native bindings may be required on some platforms)
        try:
            import textract
            text_bytes = textract.process(tmp_path)
            text = text_bytes.decode('utf-8', errors='ignore')
            # Add page count estimation for .doc files
            word_count = len(text.split())
            estimated_pages = max(1, round(word_count / 450))  # ~450 words per page
            metadata['total_pages'] = estimated_pages
            return text, metadata
        except Exception:
            pass

        # If we reach here, we couldn't process .doc; provide a clear error
        raise ValueError("Legacy .doc files require system conversion tools (pandoc/pypandoc) or the 'textract' library to be installed on the server.")

    finally:
        try:
            os.unlink(tmp_path)
        except Exception:
            pass


def analyze_doc_content(file_content: bytes, document_id: str) -> DocAnalysis:
    """
    Lightweight analysis for legacy .doc files.
    Attempts to extract plain text and derives simple structural metrics.
    """
    try:
        # Attempt to get text via the best-effort extractor
        try:
            text, _ = extract_text_from_doc(file_content, document_id)
        except Exception as e:
            logging.warning(f"[DOC_ANALYZE] Failed to extract text from .doc for analysis: {e}")
            text = ""

        # Basic paragraph detection
        lines = [l.strip() for l in text.splitlines() if l.strip()]
        total_paragraphs = len(lines)

        # Simple table heuristic: lines with tabs or pipe separators
        table_like_lines = [l for l in lines if '\t' in l or '|' in l]
        # Rough table count: group contiguous table-like lines into one table
        table_count = 0
        in_table = False
        for l in lines:
            if '\t' in l or '|' in l:
                if not in_table:
                    table_count += 1
                    in_table = True
            else:
                in_table = False

        # Images: legacy .doc extraction rarely preserves images in text extraction
        images = 0

        # Word-based page estimate (approx 450 words per page)
        total_words = sum(len(p.split()) for p in lines)
        estimated_pages = max(1, int(total_words / 450) if total_words > 0 else 1)

        # Decide strategy
        if table_count >= max(1, estimated_pages // 2):
            strategy = "hybrid"
            processing_estimate = estimated_pages * 0.8
        else:
            strategy = "text_primary"
            processing_estimate = estimated_pages * 0.2

        has_mixed = (images > 0 and total_paragraphs > 0)

        return DocAnalysis(
            total_paragraphs=total_paragraphs,
            tables=table_count,
            images=images,
            estimated_pages=estimated_pages,
            recommended_strategy=strategy,
            processing_estimate=processing_estimate,
            has_mixed_content=has_mixed
        )

    except Exception as e:
        logging.warning(f"[DOC_ANALYZE] Analysis failed for {document_id}: {e}")
        return DocAnalysis(
            total_paragraphs=0,
            tables=0,
            images=0,
            estimated_pages=1,
            recommended_strategy="text_primary",
            processing_estimate=1.0,
            has_mixed_content=False
        )

# =========================== Main Extraction Router ===========================

def extract_text_by_file_type(file_content: bytes, filename: str, file_extension: str) -> Tuple[str, Dict]:
    """
    Route file to appropriate extraction function based on file type
    Returns tuple of (extracted_text, metadata)
    """
    file_ext = file_extension.lower().lstrip('.')
    
    try:
        if file_ext == 'docx':
            return extract_text_from_docx(file_content, filename)
        elif file_ext == 'doc':
            return extract_text_from_doc(file_content, filename)
        elif file_ext in ['xlsx', 'xls']:
            return extract_text_from_excel(file_content, filename)
        elif file_ext == 'pptx':
            return extract_text_from_pptx(file_content, filename)
        elif file_ext in ['txt', 'md']:
            return extract_text_from_txt(file_content, filename)
        elif file_ext == 'pdf':
            return extract_text_from_pdf_direct(file_content, filename)
        elif file_ext in ['gdoc', 'gsheet', 'gslides']:
            return extract_text_from_google_file(file_content, filename, file_extension)
        elif file_ext in ['html', 'htm']:
            return extract_text_from_html(file_content, filename)
        elif file_ext == 'json':
            return extract_text_from_json(file_content, filename)
        elif file_ext == 'csv':
            return extract_csv_markdown(file_content, filename)
        else:
            raise ValueError(f"Unsupported file type: {file_ext}")
            
    except Exception as e:
        logging.error(f"[EXTRACTION] Failed to extract text from {filename}: {str(e)}")
        raise

# =========================== Utility Functions ===========================

def get_supported_file_types() -> Dict[str, Dict]:
    """
    Get information about supported file types and their requirements
    """
    file_types = {
        'pdf': {
            'extensions': ['.pdf'],
            'description': 'PDF documents with text extraction',
            'library': 'PyMuPDF (fitz)',
            'available': PYMUPDF_AVAILABLE
        },
        'docx': {
            'extensions': ['.docx'],
            'description': 'Microsoft Word documents',
            'library': 'python-docx',
            'available': DOCX_AVAILABLE
        },
        'excel': {
            'extensions': ['.xlsx', '.xls'],
            'description': 'Microsoft Excel spreadsheets',
            'library': 'pandas',
            'available': PANDAS_AVAILABLE
        },
        'csv': {
            'extensions': ['.csv'],
            'description': 'CSV data files',
            'library': 'pandas',
            'available': PANDAS_AVAILABLE
        },
        'pptx': {
            'extensions': ['.pptx'],
            'description': 'Microsoft PowerPoint presentations',
            'library': 'python-pptx',
            'available': PPTX_AVAILABLE
        },
        'text': {
            'extensions': ['.txt', '.md'],
            'description': 'Plain text and Markdown files',
            'library': 'built-in',
            'available': True
        },
        'html': {
            'extensions': ['.html', '.htm'],
            'description': 'HTML web pages with intelligent content extraction',
            'library': 'beautifulsoup4',
            'available': BS4_AVAILABLE
        },
        'json': {
            'extensions': ['.json'],
            'description': 'JSON data files with flattened text extraction',
            'library': 'built-in (json)',
            'available': True
        },
        'google_docs': {
            'extensions': ['.gdoc', '.docx'],
            'description': 'Google Docs (exported as Word format)',
            'library': 'google-api-python-client + python-docx',
            'available': GOOGLE_API_AVAILABLE and DOCX_AVAILABLE
        },
        'google_sheets': {
            'extensions': ['.gsheet', '.xlsx'],
            'description': 'Google Sheets (exported as Excel format)',
            'library': 'google-api-python-client + pandas',
            'available': GOOGLE_API_AVAILABLE and PANDAS_AVAILABLE
        },
        'google_slides': {
            'extensions': ['.gslides', '.pptx'],
            'description': 'Google Slides (exported as PowerPoint format)',
            'library': 'google-api-python-client + python-pptx',
            'available': GOOGLE_API_AVAILABLE and PPTX_AVAILABLE
        }
    }
    
    return file_types

def check_library_availability() -> Dict[str, bool]:
    """
    Check which text extraction libraries are available
    """
    return {
        'python-docx': DOCX_AVAILABLE,
        'pandas': PANDAS_AVAILABLE,
        'python-pptx': PPTX_AVAILABLE,
        'PyMuPDF': PYMUPDF_AVAILABLE,
        'google-api': GOOGLE_API_AVAILABLE,
        'beautifulsoup4': BS4_AVAILABLE
    }
